#!/usr/bin/env python3
"""Fit Check deck (python-pptx) — Celtics-styled slide show of the trade audit.

Same design system as the Duran strategy deck (eyebrow + Cambria titles,
white rounded cards, dot motif, dark bookend slides), in Celtics colors.
Embeds pre-rendered charts from outputs/figures/.
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches as In, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
FIG = ROOT / "outputs" / "figures"
OUT = ROOT / "outputs" / "Breaking_Down_The_Jaylen_Brown_Trade.pptx"

# palette — Celtics
GREEN = "007A33"; DARK = "0A2A1A"; GOLD = "BA9653"; INK = "20232A"
MUTED = "6B7280"; BG = "F4F5F7"; WHITE = "FFFFFF"; LINE = "D7DAE0"
RED = "B3403A"  # negative accent only
HEAD = "Cambria"; BODY = "Calibri"

W, H = 13.333, 7.5
ASPECT = {"diet": 2.279, "onoff": 1.559, "eff": 2.424, "picks": 2.418,
          "solo": 2.431, "identity": 2.76, "redis": 2.48}
FILES = {"diet": "case_for_moving_on.png",
         "onoff": "with_without_net_2025-26.png",
         "eff": "efficiency_comps.png",
         "picks": "pick_value.png",
         "solo": "tatum_first_option.png",
         "identity": "three_point_identity.png",
         "redis": "wing_redistribution.png"}


def C(hexs):
    return RGBColor.from_string(hexs)


def txt(slide, x, y, w, h, runs, size=16, color=INK, bold=False, font=BODY,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False,
        space_after=4, line_spacing=1.0, wrap=True):
    tb = slide.shapes.add_textbox(In(x), In(y), In(w), In(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    if isinstance(runs, str):
        runs = [[(runs, {})]]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        if isinstance(para, str):
            para = [(para, {})]
        for t, o in para:
            r = p.add_run(); r.text = t
            f = r.font
            f.size = Pt(o.get("size", size)); f.bold = o.get("bold", bold)
            f.italic = o.get("italic", italic); f.name = o.get("font", font)
            f.color.rgb = C(o.get("color", color))
    return tb


def card(slide, x, y, w, h, fill=WHITE, line=LINE, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                In(x), In(y), In(w), In(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill)
    if line:
        sh.line.color.rgb = C(line); sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def dot(slide, x, y, d, fill=GREEN, text=None, tcolor=WHITE, size=15):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, In(x), In(y), In(d), In(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = C(fill); sh.line.fill.background()
    sh.shadow.inherit = False
    if text is not None:
        tf = sh.text_frame; tf.word_wrap = False
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = True; r.font.name = HEAD
        r.font.color.rgb = C(tcolor)
    return sh


def image(slide, key, x, y, max_w, max_h):
    a = ASPECT[key]
    w = max_w; hh = w / a
    if hh > max_h:
        hh = max_h; w = hh * a
    ix = x + (max_w - w) / 2
    slide.shapes.add_picture(str(FIG / FILES[key]), In(ix), In(y), In(w), In(hh))
    return ix, y, w, hh


def eyebrow(slide, text, x=0.62, y=0.42, color=GREEN):
    txt(slide, x, y, 11, 0.3, [[(text.upper(), {"color": color, "bold": True,
        "size": 12.5, "font": BODY})]], space_after=0)


def title(slide, text, x=0.62, y=0.72, w=12.1, size=30, color=DARK):
    txt(slide, x, y, w, 0.95, [[(text, {"size": size, "bold": True,
        "font": HEAD, "color": color})]], space_after=0, line_spacing=1.0)


def footer(slide, dark=False):
    c = "9AB2A4" if dark else MUTED
    txt(slide, 0.62, 7.12, 12.1, 0.3,
        [[("Fit Check · Jaylen Brown Trade Audit", {"color": c, "size": 9}),
          ("      Data: nba_api / Basketball-Reference · 2024-25 & 2025-26 · July 2026",
           {"color": c, "size": 9})]], space_after=0)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(color)


prs = Presentation()
prs.slide_width = In(W); prs.slide_height = In(H)
blank = prs.slide_layouts[6]


def new(color=BG):
    s = prs.slides.add_slide(blank); bg(s, color); return s


# ---------------------------------------------------------------- 1 TITLE
s = new(DARK)
dot(s, 0.62, 0.62, 0.34, fill=GOLD)
txt(s, 1.05, 0.6, 10, 0.4, [[("BOSTON CELTICS  ·  FIT CHECK",
    {"color": "BCD3C6", "bold": True, "size": 13})]], space_after=0)
txt(s, 0.62, 2.15, 12.1, 2.0, [
    [("Breaking Down", {"size": 46, "bold": True, "font": HEAD,
      "color": WHITE})],
    [("the Jaylen Brown Trade", {"size": 46, "bold": True, "font": HEAD,
      "color": GOLD})],
], space_after=2, line_spacing=1.0)
txt(s, 0.62, 4.35, 11.6, 1.0, [[(
    "On July 1, 2026 Boston traded Brown to Philadelphia for Paul George and "
    "picks. This is a breakdown on why the trade might not be all bad news "
    "for Boston. In the world of \"that boy tuff\" and analytical basketball, "
    "this is a Celtics fan's version of this trade.",
    {"size": 17, "color": "CFE0D6"})]], line_spacing=1.12)
txt(s, 0.62, 6.6, 12, 0.4, [[("Numbers, not takes · nba_api + Basketball-"
    "Reference · unit-tested pipeline", {"color": "8FA89B", "size": 12,
    "italic": True})]], space_after=0)

# ---------------------------------------------------------------- 2 SETUP
s = new()
eyebrow(s, "The setup")
title(s, "A $53M wing that shoots long twos, and a live or die by the "
         "three pointer system", size=27)
stats = [("34.4%", "of the salary cap", "$53.1M in 2025-26 — over a third of "
          "the sheet on one wing"),
         ("+16%", "worse shot diet, YoY", "bad-shot index 0.328 → 0.379 — "
          "drifting away from the system"),
         ("4", "assets back from PHI", "Paul George, a 2028 first, a 2031 "
          "unprotected first, two seconds")]
cw, gap = 3.83, 0.31
for i, (big, lab, sub) in enumerate(stats):
    x = 0.62 + i * (cw + gap)
    card(s, x, 2.05, cw, 2.35)
    txt(s, x + 0.28, 2.28, cw - 0.5, 0.9, [[(big, {"size": 44, "bold": True,
        "font": HEAD, "color": RED if i == 1 else GREEN})]], space_after=0)
    txt(s, x + 0.3, 3.28, cw - 0.55, 0.4, [[(lab, {"size": 15, "bold": True,
        "color": DARK})]], space_after=0)
    txt(s, x + 0.3, 3.66, cw - 0.55, 0.7, [[(sub, {"size": 12.5,
        "color": MUTED})]], line_spacing=1.05)
txt(s, 0.62, 4.75, 12.1, 1.4, [
    [("Boston's edge is shooting more threes than opponents. ",
      {"bold": True, "color": INK, "size": 15}),
     ("The question this deck answers with data: did Brown's game pull with "
      "that identity or against it, and was the return worth it?",
      {"color": INK, "size": 15})],
], line_spacing=1.2)
footer(s)

# ---------------------------------------------------------------- 3 SHOT DIET
s = new()
eyebrow(s, "The shot diet")
title(s, "The shot profile was drifting away from what Boston wins with")
image(s, "diet", 0.5, 1.75, 7.6, 3.5)
tx = 8.45
rows = [("Fewer threes, doubled long twos", RED,
         "3PT rate fell 0.32 → 0.26; long-two rate doubled 0.07 → 0.14 — the "
         "one shot the system exists to kill."),
        ("Two-thirds self-created", DARK,
         "Iso / 3+ dribble rate rose 0.53 → 0.64, with more late-clock "
         "bailouts (0.18 → 0.22)."),
        ("No tough-shot premium", GREEN,
         "The xFG model (2,661 shots) scores his shot-making at expectation "
         "— ≈ 0 both years. The difficulty is self-inflicted.")]
yy = 1.95
for head, col, body in rows:
    dot(s, tx, yy + 0.04, 0.16, fill=col)
    txt(s, tx + 0.32, yy, 4.3, 1.2, [[(head, {"bold": True, "size": 14,
        "color": DARK})], [(body, {"size": 12.3, "color": MUTED})]],
        line_spacing=1.06, space_after=2)
    yy += 1.32
txt(s, 0.62, 6.1, 12.1, 0.75, [[("Brown's own shot selection drifted far "
    "away from what the system is built on, ", {"size": 14, "color": INK}),
    ("and his \"tough\" shot making isn't good enough to excuse it.",
     {"size": 14, "bold": True, "color": DARK})]], line_spacing=1.1)
footer(s)

# ---------------------------------------------------------------- 4 GAME PLAN
s = new()
eyebrow(s, "The game plan")
title(s, "Boston's identity is math and Brown ruins it")
image(s, "identity", 0.7, 1.7, 11.9, 4.35)
txt(s, 0.62, 6.05, 12.1, 1.0, [
    [("#1 · #1 · #4 in three-point volume and a top-2 offense all three "
      "years. Threes returned 1.10 points per shot vs 0.97 for twos. ",
      {"size": 13, "color": INK}),
     ("Brown: last of 7 perimeter players in 3PT rate (.263 vs team .467), "
      "first in long twos. ", {"size": 13, "bold": True, "color": RED}),
     ("A shot-mix problem, not a finishing one (his 1.046 PPS vs 1.101).",
      {"size": 13, "color": MUTED})],
], line_spacing=1.15)
footer(s)

# ---------------------------------------------------------------- 5 SHOTS GO
s = new()
eyebrow(s, "Where the shots go")
title(s, "Redistribute Brown's 21.7 shots to the wings")
image(s, "redis", 0.9, 1.7, 11.5, 4.3)
txt(s, 0.62, 6.0, 12.1, 1.05, [
    [("Ceiling: +1.3 pts/g (~3.5 wins) if every wing's efficiency survived "
      "the extra volume — it won't. ", {"size": 13, "bold": True,
      "color": GREEN}),
     ("Usage-adjusted: roughly a wash (−0.7 to +0.9 wins). And the George-"
      "only 1-for-1 is negative for the Celtics. ",
      {"size": 13, "color": INK}),
     ("The real upside is Tatum's creation (5.7-6.1 ast, .498 3PT-rate "
      "gravity) upgrading the shots these wings get, and that effect is "
      "already measured in the +11.9 Tatum-led lineups. Combined honest "
      "range: ~+4 to +8 wins (without Brown this year).", {"size": 13,
      "bold": True, "color": DARK})],
], line_spacing=1.15)
footer(s)

# ---------------------------------------------------------------- 6 CONTINGENT
s = new()
eyebrow(s, "Contingent value")
title(s, "The production needed Tatum and White to prop it up")
image(s, "onoff", 0.5, 1.7, 6.6, 4.6)
tx = 7.5
rows = [("With Tatum: +16.2 · without: +4.8", DARK,
         "An 11.3-point swing in 2025-26. Brown's lineups won when the "
         "infrastructure was on the floor."),
        ("With White: +9.4 · without: −1.0", DARK,
         "Same story with the other connector — Brown-led lineups without "
         "White were underwater."),
        ("The honest caveat", GOLD,
         "Tatum's injury contaminates the 2025-26 split; the confound is "
         "documented, not hidden. It softens — but doesn't reverse — the "
         "pattern.")]
yy = 1.95
for head, col, body in rows:
    dot(s, tx, yy + 0.04, 0.16, fill=col)
    txt(s, tx + 0.32, yy, 5.2, 1.3, [[(head, {"bold": True, "size": 14.5,
        "color": DARK})], [(body, {"size": 12.8, "color": MUTED})]],
        line_spacing=1.08, space_after=2)
    yy += 1.5
txt(s, tx, yy + 0.05, 5.2, 0.6, [[("A $53M player whose value is contingent "
    "on his co-stars is a fit problem, not a star.", {"italic": True,
    "bold": True, "size": 13, "color": GREEN})]], line_spacing=1.05)
footer(s)

# ---------------------------------------------------------------- 5 HIERARCHY
s = new()
eyebrow(s, "The hierarchy test")
title(s, "First-option Tatum will be a problem (in a good way)")
image(s, "solo", 0.5, 1.7, 7.4, 3.2)
tx = 8.3
rows5 = [("25-4 (86.2%) when Brown sat", GREEN,
         "29 games, 2023-present, +13.3/game — including 10-0 in the "
         "title year. Both healthy: 73.7%. Lineups: Tatum-led +11.9 > "
         "together +7.6 > neither +3.5 (control)."),
        ("Usage up, efficiency mostly holds", DARK,
         "Pooled: 28.0/7.9/5.7 at 34.6% usage (+4.7 pts of usage), TS "
         ".593 → .572 — a modest efficiency dip, printed, not hidden."),
        ("Worth ~5-7 wins a season", GOLD,
         "Projection, not observation: the +4.3 lineup gap, shrunk 40-60% "
         "for bench/garbage contamination → +1.7-2.6 team net → +4.7-7.0 "
         "wins.")]
yy = 1.9
for head, col, body in rows5:
    dot(s, tx, yy + 0.04, 0.16, fill=col)
    txt(s, tx + 0.32, yy, 4.45, 1.25, [[(head, {"bold": True, "size": 14,
        "color": DARK})], [(body, {"size": 12.2, "color": MUTED})]],
        line_spacing=1.06, space_after=2)
    yy += 1.44
txt(s, 0.62, 5.05, 7.4, 1.3, [[("Tatum's projected solo season: ",
    {"size": 13, "color": INK}),
    ("27.2-28.8 pts · 7.8-8.9 reb · 5.1-6.2 ast · .566-.584 TS",
     {"size": 13, "bold": True, "color": GREEN})]], line_spacing=1.12)
footer(s)

# ---------------------------------------------------------------- 6 CONTRACT
s = new()
eyebrow(s, "The contract")
title(s, "Max-contract price, below-median-efficiency production")
image(s, "eff", 1.2, 1.65, 10.9, 4.55)
txt(s, 0.62, 6.35, 12.1, 0.7, [[("6.9 WS on $53.1M = $7.7M per win share at "
    "34% of the cap, with a 0.573 TS%, below the max-wing median. ",
    {"size": 14, "color": INK}),
    ("Picks 15-30 buy the same wins at ~$2.1M/WS on rookie deals.",
     {"size": 14, "bold": True, "color": GREEN})]], line_spacing=1.1)
footer(s)

# ---------------------------------------------------------------- 6 RETURN
s = new()
eyebrow(s, "The return")
title(s, "George is the cleaner fit for the Celtics system, and new lottery odds", size=28)
comp = [("3PT rate", "0.262", "0.497"),
        ("Catch-&-shoot rate", "0.163", "0.409"),
        ("Iso / 3+ dribble rate", "0.639", "0.404"),
        ("Contested-shot rate", "0.517", "0.370"),
        ("Bad-shot index", "0.379", "0.273")]
card(s, 0.62, 1.95, 7.2, 4.35)
txt(s, 0.94, 2.18, 6.6, 0.4, [[("2025-26 shot profile", {"size": 15,
    "bold": True, "font": HEAD, "color": DARK})]], space_after=0)
txt(s, 3.9, 2.62, 1.9, 0.3, [[("BROWN", {"size": 11.5, "bold": True,
    "color": MUTED})]], space_after=0)
txt(s, 5.9, 2.62, 1.9, 0.3, [[("GEORGE", {"size": 11.5, "bold": True,
    "color": GREEN})]], space_after=0)
yy = 3.0
for lab, jb, pg in comp:
    txt(s, 0.94, yy, 2.9, 0.35, [[(lab, {"size": 13, "color": INK})]],
        space_after=0)
    txt(s, 3.9, yy, 1.9, 0.35, [[(jb, {"size": 13.5, "bold": True,
        "color": MUTED})]], space_after=0)
    txt(s, 5.9, yy, 1.9, 0.35, [[(pg, {"size": 13.5, "bold": True,
        "color": GREEN})]], space_after=0)
    yy += 0.6
txt(s, 0.94, 5.62, 6.6, 0.5, [[("George takes the shots the system is "
    "built to generate. He plugs in instead of stopping it.",
    {"size": 11.5, "italic": True, "color": MUTED})]], line_spacing=1.05)
card(s, 8.15, 1.95, 4.55, 4.35, fill=DARK, line=None)
txt(s, 8.45, 2.2, 4.0, 0.4, [[("PLUS THE PICKS", {"size": 12.5, "bold": True,
    "color": "BCD3C6"})]], space_after=0)
txt(s, 8.45, 2.7, 4.0, 1.0, [
    [("2028 first + 2031", {"size": 20, "bold": True, "font": HEAD,
      "color": WHITE})],
    [("unprotected first", {"size": 20, "bold": True, "font": HEAD,
      "color": WHITE})]], space_after=2, line_spacing=1.05)
txt(s, 8.45, 3.85, 4.0, 1.6, [
    [("Under the post-2019 flattened lottery exact odds enumerated across "
      "all 24,024 seed permutations, late-lottery firsts convey materially "
      "more top-4 equity than the old system.", {"size": 12.5,
      "color": "CFE0D6"})]], line_spacing=1.15)
txt(s, 8.45, 5.55, 4.0, 0.6, [[("Cheap, controllable, tradeable — the assets "
    "a capped-out roster can't otherwise get.", {"size": 11.5,
    "italic": True, "color": "9AB2A4"})]], line_spacing=1.08)
footer(s)

# ---------------------------------------------------------------- 7 WEAKEST
s = new()
eyebrow(s, "Where the case is weakest")
title(s, "The counter-evidence, on the record")
weak = [("Brown could carry", GOLD,
         "His 7-game solo cell was the best in the 2024-25 matrix (6-1, "
         "+15.8), and he carried 58 games alone in 2025-26 at 28.9 ppg "
         "(36-22, +5.7). The man can be a first option — just not here."),
        ("The defense is a real loss", GOLD,
         "Brown took 19% of his defended possessions against opponents' #1 "
         "options (Tatum: 8%) at better points-allowed, with -4.2/-5.7 "
         "defended-FG% margins. Boston traded its toughest assignment."),
        ("The turnover tax wasn't his", RED,
         "Play-by-play says Brown's live-ball TO rate (1.75/36) was LOWER "
         "than Tatum's (1.84) in 2024-25. That angle exonerates him — and "
         "the media consensus graded the trade for Philadelphia.")]
cw, gap = 3.83, 0.31
for i, (head, col, body) in enumerate(weak):
    x = 0.62 + i * (cw + gap)
    card(s, x, 2.05, cw, 3.3)
    dot(s, x + 0.3, 2.33, 0.42, fill=col, text=str(i + 1), size=16)
    txt(s, x + 0.3, 2.97, cw - 0.55, 0.5, [[(head, {"size": 16, "bold": True,
        "font": HEAD, "color": DARK})]], space_after=0)
    txt(s, x + 0.3, 3.5, cw - 0.55, 1.7, [[(body, {"size": 12.8,
        "color": MUTED})]], line_spacing=1.12)
txt(s, 0.62, 5.7, 12.1, 0.9, [[("Brown as an individual is a remarkable "
    "basketball player, there is no denying that. ", {"size": 14,
    "bold": True, "color": DARK}), ("But he was just an expensive player "
    "that didn't quite fit what the Celtics are trying to do. The data says "
    "the best version of this team was always Tatum at the top.",
    {"size": 14, "color": MUTED})]], line_spacing=1.1)
footer(s)

# ---------------------------------------------------------------- 8 CLOSER
s = new(DARK)
dot(s, 0.62, 0.7, 0.34, fill=GOLD)
txt(s, 1.05, 0.68, 10, 0.4, [[("BOTTOM LINE", {"color": "BCD3C6",
    "bold": True, "size": 13})]], space_after=0)
lines = [("The fit was real and it was drifting.", "Fewer threes, doubled "
          "long twos, two-thirds self-created — with no shot-making premium "
          "to pay for it."),
         ("The value was contingent.", "Brown-led lineups needed Tatum and "
          "White on the floor; at 34% of the cap, contingent is expensive."),
         ("The return buys identity and optionality.", "George takes the "
          "system's shots; two firsts add the cheap, controllable assets a "
          "capped roster can't otherwise acquire.")]
yy = 1.95
for head, body in lines:
    dot(s, 0.7, yy + 0.06, 0.18, fill=GOLD)
    txt(s, 1.05, yy, 11.4, 1.0, [[(head + "  ", {"bold": True, "size": 21,
        "font": HEAD, "color": WHITE}), (body, {"size": 15,
        "color": "CFE0D6"})]], line_spacing=1.1, space_after=2)
    yy += 1.35

prs.save(str(OUT))
print(f"wrote {OUT} ({len(prs.slides._sldIdLst)} slides)")
